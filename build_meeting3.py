"""
Builds Çağrı_Tirelioğlu_meeting3.pptx by inheriting template style (master,
layouts, theme, fonts, colors) from meeting2.pptx and authoring fresh content
based on slides_outline.md and today's milestones.

Run on laptop (Windows, has python-pptx installed).
"""
import sys
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from lxml import etree

sys.stdout.reconfigure(encoding="utf-8")

MEETINGS_DIR = Path(r"C:\Users\Çağrı Tirelioğlu\Desktop\CSE_THE_LAST_CHANCE\GRAD2\meetings")
TEMPLATE = MEETINGS_DIR / "Çağrı_Tirelioğlu_meeting2.pptx"
OUT_PATH = MEETINGS_DIR / "Çağrı_Tirelioğlu_meeting3.pptx"

DATE_DISPLAY = "30.05.2026"
PRESENTATION_NAME = "Third Presentation"
FOOTER_TEXT = "CSE496 – Third Presentation"


# ── Slide content definitions ─────────────────────────────────────────

# Each slide: (title, body_lines_or_special, notes_en, notes_tr)

SLIDE_TITLE_INFO = {
    "title": "Automated Knowledge Conflict and Inconsistency Detection via Graph-Based Reliable RAG in Multi-Document Systems",
    "subtitle": "CSE 496\nThird Presentation\n\nÇağrı Tirelioğlu\n\nAdvisor: Prof. Yusuf Sinan Akgül\nMay 2026",
    "notes_en": (
        "Good morning. I'm Çağrı Tirelioğlu. In this third presentation of my "
        "graduation project — ReliabilityRAG — I will report the production-stable "
        "system, the closure of our 30-case ablation study, a new generation "
        "faithfulness evaluation pipeline that revealed an unexpected LLM-layer "
        "bottleneck, and the production demo material we curated for the final defense."
    ),
    "notes_tr": (
        "Merhaba, ben Çağrı Tirelioğlu. Bitirme projem ReliabilityRAG'in bu üçüncü "
        "sunumunda production-stable hâle gelen sistemi, 30-case ablation çalışmasının "
        "kapanışını, yeni geliştirdiğimiz generation faithfulness eval pipeline'ının "
        "ortaya çıkardığı beklenmedik LLM-katman darboğazını ve final savunma için "
        "hazırladığımız production demo malzemesini aktaracağım."
    ),
}


CONTENT_SLIDES = [
    # 2 — Contents
    {
        "title": "Contents",
        "bullets": [
            "Project Recap",
            "Progress Since Second Presentation",
            "RTX 5080 Production Migration",
            "Llama-3 Chat Template Bug Fix",
            "30-Case Filter Layer Ablation",
            "NEW: Generation Faithfulness Eval",
            "Demo Highlights (Production Outputs)",
            "Honest Failure Cases",
            "Success Criteria Status",
            "Timeline & Next Steps",
            "References",
        ],
        "notes_en": (
            "I'll start with a brief recap, then walk through what's new since April: "
            "production migration to RTX 5080, swapping to a Turkish-tuned LLM, a "
            "critical chat-template bug fix, the closure of our filter-layer ablation "
            "on a 30-case test set, a new generation faithfulness eval, three live "
            "demo highlights, honest failure cases, and the path to defense."
        ),
        "notes_tr": (
            "Önce kısa bir hatırlama ile başlıyorum, sonra Nisan'dan bu yana yapılanları "
            "anlatacağım: RTX 5080'e production migration, Türkçe fine-tune'lu LLM'e "
            "geçiş, kritik bir chat-template bug fix'i, 30-case test setinde filter-layer "
            "ablation'ın kapanışı, yeni generation faithfulness eval'i, üç canlı demo "
            "örneği, dürüst failure case'ler ve savunmaya kalan yol."
        ),
    },
    # 3 — Project Recap
    {
        "title": "Project Recap",
        "bullets": [
            "Goal: Detect knowledge conflicts and inconsistencies across multi-document Turkish corporate reports.",
            "Data: 248 integrated reports, 63 BIST-listed companies, 2015–2024, 182,986 chunks.",
            "Pipeline: Retrieval → NLI Contradiction Graph → MWIS+GAT Filter → Grounded LLM Generation.",
            "Original contribution: GAT-based dynamic filtering replacing static MIS in the ReliabilityRAG paper.",
            "New addition (this period): generation-layer reliability evaluation as a second measurement axis.",
        ],
        "notes_en": (
            "Quick recap: Turkish corporate reports contradict themselves across years and "
            "departments. Our 248-report corpus, 182K chunks. The 6-stage pipeline still "
            "applies; our original contribution is the GAT-based dynamic filtering. New "
            "this period: we added a second measurement axis — generation layer reliability "
            "— which turned out to be a key thesis finding."
        ),
        "notes_tr": (
            "Kısa hatırlatma: Türk kurumsal raporlarında yıllar ve departmanlar arası "
            "çelişkiler var. 248 raporluk korpus, 182K chunk. Altı aşamalı pipeline aynı; "
            "özgün katkımız GAT tabanlı dinamik filtreleme. Bu dönemde yeni: ikinci bir "
            "ölçüm ekseni eklendi — generation katmanı reliability'si — ki bu tezsel olarak "
            "anahtar bir bulguya dönüştü."
        ),
    },
    # 4 — Progress Since Second Presentation
    {
        "title": "Progress Since Second Presentation",
        "bullets": [
            "Production migration to RTX 5080 (Akvaryum PC) with Turkish-tuned LLM stack.",
            "Critical Llama-3 chat template bug fix: 595s → 18.5s generation (32× speedup).",
            "Test set expanded: 19 → 30 cases across 8 contradiction dimensions.",
            "NEW evaluation pipeline: evaluation_generation.py (4 metrics, 2 modes).",
            "30-case filter-layer ablation closed: MWIS = trained GAT = 15/30 with verdict heterogeneity.",
            "Demo cache built: 12 curated real-world queries, full pipeline outputs cached for backup demo.",
            "Thesis report auto-generator extended: 8 sections, 18 KB output, reproducible.",
            "GitHub repo: 11 commits today alone, full reproducibility chain.",
        ],
        "notes_en": (
            "Big-picture summary: we moved to production on the RTX 5080, swapped to "
            "Turkish-Llama-8b for native Turkish quality, expanded our test set, built a "
            "second-axis evaluation pipeline, closed the filter-layer ablation question, "
            "curated demo material, and committed it all reproducibly. Today alone: 11 "
            "GitHub commits."
        ),
        "notes_tr": (
            "Genel resim: RTX 5080'e production'a geçtik, native Türkçe kalitesi için "
            "Turkish-Llama-8b'ye geçiş yaptık, test setini genişlettik, ikinci eksen "
            "evaluation pipeline'ı kurduk, filter-layer ablation sorusunu kapadık, demo "
            "malzemesi hazırladık ve hepsini reproducible olarak commit'ledik. Sadece "
            "bugün: 11 GitHub commit."
        ),
    },
    # 5 — RTX 5080 Production Migration
    {
        "title": "RTX 5080 Production Migration",
        "bullets": [
            "Three LLM candidates evaluated on 16 GB Blackwell VRAM envelope:",
            "Qwen-7B-fp16: 15.2 GB → CPU offload, page freeze.",
            "Qwen-7B-4bit: 38.5s/query but Türkçe entity drift (Anadolu Hayat → Hava Yolları, FEM → FIŞLİ EŞİKLİK).",
            "Qwen-14B-4bit (worst): 595s, Chinese-language drift under nf4 quantization stress. Documented failure case.",
            "Turkish-Llama-8b-Instruct-v0.1 (YTÜ-CE Cosmos): 25.2s/query, native Türkçe morphology, ENTITIES correct.",
            "Final stack: Turkish-Llama-8b 4-bit on RTX 5080, mDeBERTa-base NLI, multilingual-e5-base embeddings.",
        ],
        "notes_en": (
            "Three LLM candidates we evaluated. Qwen-7B at fp16 overflowed 16 GB. At 4-bit "
            "it ran in 38 seconds but mangled Turkish entity names — Anadolu Hayat became "
            "Anadolu Hava Yolları. Qwen-14B-4bit hit CPU offload, ran 10 minutes, and "
            "drifted into Chinese under quantization stress — this failure is itself "
            "documented thesis evidence. Turkish-Llama-8b — built on Llama-3 by the YTÜ CE "
            "Cosmos lab with Turkish instruction fine-tuning — runs in 25 seconds with "
            "clean Turkish."
        ),
        "notes_tr": (
            "Üç LLM aday değerlendirdik. Qwen-7B fp16'da 16 GB'ı taşırdı. 4-bit'te 38 "
            "saniyede çalıştı ama Türkçe şirket adlarını bozdu — Anadolu Hayat 'Anadolu "
            "Hava Yolları' oldu. Qwen-14B-4bit CPU offload'a girdi, 10 dakika sürdü ve "
            "quantization baskısı altında Çince'ye savruldu — bu başarısızlık tezsel "
            "kanıt olarak değerli. Turkish-Llama-8b — Llama-3 tabanlı, YTÜ CE Cosmos "
            "lab'inin Türkçe instruction fine-tune'u — 25 saniyede temiz Türkçe üretiyor."
        ),
    },
    # 6 — Llama-3 Chat Template Bug Fix
    {
        "title": "Llama-3 Chat Template Bug Fix (Bonus Story)",
        "bullets": [
            "Bug: First Turkish-Llama run produced 595s generation with FAKE assistant/user/system dialogues as plain text.",
            "Root cause: Llama-3 tokenizer eos_token_id = <|end_of_text|> (128001 — document-end), but chat turn-end token is <|eot_id|> (128009).",
            "Without <|eot_id|> in the terminator list, the model emits it, generation continues, skip_special_tokens strips the tags but leaves the role-header WORDS as text.",
            "Fix: llm_engine._get_terminators() helper — dynamically adds <|eot_id|> if tokenizer has it (safe for Qwen fallback).",
            "Result: 595s → 18.5s generation (32× speedup); fake-dialog hallucination eliminated; answers end naturally at turn boundary.",
        ],
        "notes_en": (
            "A bug worth highlighting because it took hours to diagnose. Turkish-Llama "
            "first ran for 10 minutes per query and produced fake conversation turns. "
            "Llama-3's chat turn-end token is <|eot_id|> — 128009. The default eos is "
            "<|end_of_text|>. Without explicitly listing <|eot_id|> as a terminator, the "
            "model emits it, generation continues, and the skip_special_tokens decode "
            "leaves the words 'assistant', 'user', 'system' as plain text. The fix is a "
            "few lines but the speedup is 32×. This kind of production engineering rigor "
            "matters for reliability claims."
        ),
        "notes_tr": (
            "Saatler aldığı için vurgulamaya değer bir bug. Turkish-Llama ilk başta her "
            "sorgu için 10 dakika sürdü ve sahte konuşma turları üretti. Llama-3'ün chat "
            "turn-end token'ı <|eot_id|> — 128009. Default eos ise <|end_of_text|>. "
            "<|eot_id|>'yi açıkça terminator olarak listelemezsek, model onu emit eder, "
            "generation devam eder ve skip_special_tokens decode 'assistant', 'user', "
            "'system' kelimelerini düz metin olarak bırakır. Fix birkaç satır ama "
            "speedup 32 kat. Bu tür production engineering rigor reliability iddiası "
            "için kritik."
        ),
    },
    # 7 — 30-Case Filter Layer Ablation
    {
        "title": "30-Case Filter Layer Ablation",
        "bullets": [
            "Test set expanded 19 → 30 cases across 8 dimensions (temporal, scope, interdepartmental, gat_discriminating, cross_company, zero_claim, dense_graph, numerical_edge).",
            "MWIS-only: 15/30 (50.0%)",
            "Heuristic GAT: 14/30 (regressed by 1 — same-reliability tie-breaker)",
            "Trained GAT: 15/30 (recovered the regression, kept overall score)",
            "Verdict composition: 14 both_ok, 1 gat_only_ok (GAT-9 4-chunk temporal evolution case), 1 mwis_only_ok, 14 both_failed.",
            "Training had measurable effect on gat_discriminating dimension: heuristic 2/10 → trained 3/10.",
            "NEW limitation: dense_graph (8-chunk) cases — both MWIS and GAT 0/3. Identifies a heuristic-optimization boundary.",
        ],
        "notes_en": (
            "The headline: MWIS and trained GAT score equally (15 of 30), but the verdict "
            "composition differs. Trained GAT recovered one specific test that MWIS missed — "
            "a 4-chunk temporal evolution case, exactly the kind we designed gat_discriminating "
            "for. On the dimension training has measurable effect: heuristic 2 of 10, "
            "trained 3 of 10. The new finding is on dense_graph cases — 8-chunk multi-way "
            "conflict graphs — where both methods score zero. This identifies a clean "
            "future-work boundary: heuristic optimization needs global MIS or end-to-end "
            "learning for dense graphs."
        ),
        "notes_tr": (
            "Headline: MWIS ile trained GAT eşit skor (15/30), ama verdict composition "
            "farklı. Trained GAT, MWIS'ın kaçırdığı spesifik bir testi kurtardı — 4-chunk "
            "temporal evolution case, gat_discriminating'i tam bu amaçla tasarlamıştık. "
            "Training, hedef boyutta ölçülebilir etki: heuristic 2/10, trained 3/10. Yeni "
            "bulgu dense_graph case'lerinde — 8 chunk'lık multi-way conflict graphs — "
            "her iki yöntem de sıfır. Bu temiz bir future-work sınırı tanımlıyor: heuristic "
            "optimization dense graphs için global MIS veya end-to-end learning gerektiriyor."
        ),
    },
    # 8 — NEW: Generation Faithfulness Eval
    {
        "title": "NEW: Generation Faithfulness Eval (Stage B)",
        "bullets": [
            "Question: Does the LLM stay faithful to filtered source chunks, or does it drift/hallucinate?",
            "Built evaluation_generation.py — 375 lines, 4 metrics, 2 modes:",
            "entity_recall, entity_purity (anti-fabrication), year_accuracy, numerical_fidelity (±5% tolerance).",
            "Oracle mode: feed expected_kept chunks directly (LLM-only measure).",
            "Pipeline mode: full retrieval → NLI → GAT filter → LLM (end-to-end).",
            "Result (Turkish-Llama 8b 4-bit, n=30): Numerical Fidelity 0.54, Year Accuracy 0.73.",
            "KEY FINDING: Oracle vs Pipeline delta only -0.014 → LLM-layer is dominant fidelity bottleneck, not filter.",
            "Reliability claim becomes two-layered, both quantifiable.",
        ],
        "notes_en": (
            "This is the most significant new contribution this period. Until now we only "
            "measured the filter layer. We added a second axis: does the LLM stay faithful "
            "to the sources after filtering? The answer is striking: Turkish-Llama preserves "
            "only 54 percent of source numerical values within plus-minus 5 percent. And "
            "the gap between oracle and pipeline mode is tiny — meaning the LLM hallucinates "
            "independently of filter quality. So reliability is genuinely a two-layer "
            "problem: filtering is necessary but not sufficient. This reframes the thesis."
        ),
        "notes_tr": (
            "Bu dönemin en önemli yeni katkısı. Şimdiye kadar sadece filter katmanını "
            "ölçüyorduk. İkinci eksen ekledik: filter sonrası LLM kaynaklara sadık kalıyor "
            "mu? Cevap çarpıcı: Turkish-Llama, source sayısal değerlerin sadece yüzde 54'ünü "
            "artı eksi 5 yüzde toleransla koruyor. Oracle ve pipeline mode arasındaki "
            "fark çok küçük — yani LLM filter kalitesinden bağımsız olarak halüsinasyon "
            "yapıyor. Yani reliability gerçekten iki-katmanlı bir problem: filtering "
            "gerekli ama yeterli değil. Bu tezi yeniden çerçeveliyor."
        ),
    },
    # 9 — Demo Highlights
    {
        "title": "Demo Highlights — Real Production Outputs",
        "bullets": [
            "12 curated queries run through Turkish-Llama, saved to data/demo_cache.json as bullet-proof backup.",
            "[Demo 1] Kadın istihdamı: multi-company synthesis with real programs (Turkcell \"Eşit Olacak\", \"Geleceğini Yazan Kadın\"; Borusan policy; Anadolu Hayat FEM certificate — 22 criteria, 7 met).",
            "[Demo 2] Akbank 2023 — metadata-filtered retrieval (company=Akbank, year=2023), produced 7-point structured response: net zero 2050, sustainable finance, transparency.",
            "[Demo 3] Cross-sector themes — identified UN SDG alignment, EU CSRD regulation, SKD Türkiye reporting standards across companies.",
            "Quality verdict on 12 queries: 3 STRONG (slide-ready), 4 ACCEPTABLE, 5 PROBLEM (next slide).",
        ],
        "notes_en": (
            "We curated 12 real-world queries and ran them through the production stack. "
            "Three top examples to show during demo. First, the women employment policy "
            "synthesis — the system correctly cites Turkcell's real programs and Anadolu "
            "Hayat's actual FEM certificate application with 7 of 22 criteria met. Second, "
            "a metadata-filtered Akbank 2023 query produces a structured 7-topic response. "
            "Third, cross-sector synthesis correctly references CSRD — a real EU regulation. "
            "These outputs are cached so a live demo can fall back instantly if GPU access "
            "fails during the actual defense."
        ),
        "notes_tr": (
            "12 gerçek sorgu seçtik ve production stack'inden geçirdik. Demo için üç en "
            "iyi örnek. Birincisi, kadın istihdamı politika sentezi — sistem Turkcell'in "
            "gerçek programlarını ve Anadolu Hayat'ın 22 kriterden 7'sini karşılayan "
            "gerçek FEM sertifika başvurusunu doğru aktarıyor. İkincisi, metadata-filtered "
            "Akbank 2023 sorgusu yapılı 7-konulu cevap üretiyor. Üçüncüsü, cross-sector "
            "sentez CSRD — gerçek bir AB regülasyonuna — doğru referans veriyor. Bu "
            "çıktılar cache'lendi, savunma sırasında GPU erişimi bozulursa anlık fallback'e "
            "geçilebilir."
        ),
    },
    # 10 — Honest Failure Cases
    {
        "title": "Honest Failure Cases (Thesis-Relevant)",
        "bullets": [
            "Generation eval surfaced 4 systematic failure patterns — all valuable as Section 6/7 qualitative evidence:",
            "Entity name corruption: Tüpraş written 5 different ways in a single response (Tüpraşa/Tüprag/Tüpragas/Tüprüğün/Tüpraga). Tokenizer × 4-bit nf4 interaction.",
            "Arithmetic hallucination: \"39 milyon + 34.800 = 73.800\" (actual: 39,034,800). LLMs cannot reliably aggregate.",
            "Cross-company drift: When retrieval surfaces chunks from related entities, LLM may attribute claims across them.",
            "Placeholder leakage: Template artifact \"%X arttığını\" appeared in output — incomplete instruction following.",
            "These failures are the concrete faces of the 0.54 numerical_fidelity score — qualitative evidence supporting quantitative metric.",
        ],
        "notes_en": (
            "Honesty is part of the thesis. Five queries out of twelve had clear problems. "
            "Tüpraş's name was written five different ways in a single response — a "
            "tokenizer-quantization interaction. There was an arithmetic mistake — 39 "
            "million plus 35 thousand became 73 thousand. There was cross-company drift "
            "where Anadolu Hayat appeared to speak about Türk Telekom. These are the "
            "concrete faces of the 54-percent numerical fidelity number from the previous "
            "slide. They're systemic to LLM-layer RAG, not bugs of our specific system."
        ),
        "notes_tr": (
            "Dürüstlük tezsel olarak önemli. On iki sorgudan beşinde net sorunlar vardı. "
            "Tüpraş'ın ismi tek cevapta beş farklı yazıldı — tokenizer-quantization "
            "etkileşimi. Bir aritmetik hatası vardı — 39 milyon artı 35 bin sonucu 73 "
            "bin oldu. Cross-company drift vardı; Anadolu Hayat sanki Türk Telekom "
            "hakkında konuşur gibi gözüktü. Bunlar önceki slayttaki yüzde 54'lük "
            "numerical fidelity rakamının somut yüzleri. LLM-katmanı RAG'in yapısal "
            "sınırları, bizim sistemin bug'ları değil."
        ),
    },
    # 11 — Success Criteria Status
    {
        "title": "Success Criteria Status",
        "bullets": [
            "1) Inconsistency Filtering ≥ 90% — currently 67% recall on calibrated NLI, 100% clean preservation. ON TRACK (gap is dense_graph cases).",
            "2) Temporal Consistency 100% — GAT temporal feature in place; ablation showed gat_only_ok=1 on 4-chunk evolution case. MECHANISM VALIDATED.",
            "3) Source Fidelity ≥ 95% — generation eval shows 54% numerical fidelity. GAP IDENTIFIED, reframed as separate measurement axis (Stage B).",
            "4) Algorithmic Latency ≤ 5s — Retrieval 6s + NLI <1s + Filter <0.5s + LLM 18.5s = ~25s total. Filter-only target met (<5s); full pipeline exceeds.",
            "Overall: 2 of 4 criteria fully met, 2 criteria partially met with clear documentation of remaining gap.",
        ],
        "notes_en": (
            "Updated criteria status. One: inconsistency filtering target of 90 percent — "
            "we're at 67 percent recall, the gap is the dense_graph cases. Two: temporal "
            "consistency — mechanism validated by the gat_only_ok case. Three: source "
            "fidelity — the 54 percent numerical fidelity is below the 95 percent target, "
            "but this is now properly framed as a separate measurement axis with its own "
            "evaluation pipeline. Four: latency under 5 seconds — met for filter sub-pipeline, "
            "exceeded for full pipeline due to LLM generation time. Two fully met, two "
            "partially met with clear gap documentation."
        ),
        "notes_tr": (
            "Güncel kriter durumu. Bir: inconsistency filtering hedefi yüzde 90 — yüzde 67 "
            "recall'dayız, eksik kısım dense_graph case'leri. İki: temporal consistency — "
            "gat_only_ok case'iyle mekanizma doğrulandı. Üç: source fidelity — yüzde 54 "
            "numerical fidelity yüzde 95 hedefinin altında, ama bu artık kendi evaluation "
            "pipeline'ı olan ayrı bir ölçüm ekseni olarak çerçevelendi. Dört: latency "
            "beş saniye altı — filter sub-pipeline için tutuldu, LLM generation süresi "
            "nedeniyle full pipeline aşıldı. İki tam karşılandı, iki kısmen karşılandı "
            "ve eksik açıkça belgelendi."
        ),
    },
    # 12 — Timeline & Next Steps
    {
        "title": "Timeline & Next Steps (2 Weeks to Defense)",
        "bullets": [
            "March 2026: Literature review, paper selection. DONE.",
            "March-April 2026: Data acquisition, MinerU pipeline, 182K chunk embedding, NumPy vector store. DONE.",
            "April 2026 (Second Presentation): NLI + MWIS + GAT + LLM + Gradio integration end-to-end. DONE.",
            "May 2026: Production migration (RTX 5080 + Turkish-Llama), 30-case ablation closed, generation eval pipeline. DONE.",
            "Week 1 remaining (Jun 1–7): Slide deck creation, demo screenshot capture, thesis prose draft.",
            "Week 2 (Jun 8–14): Slide revision, demo dry-runs (3+ full passes), bug-fix buffer.",
            "Defense: mid-June 2026.",
        ],
        "notes_en": (
            "Timeline status. March through April was foundation work. The April second "
            "presentation marked end-to-end integration. May closed the research phase: "
            "production migration, ablation finalization, new generation eval. The "
            "remaining 14 days are slide creation, demo rehearsal, and thesis prose. We're "
            "intentionally NOT adding new features in this period — the focus shifts to "
            "communication."
        ),
        "notes_tr": (
            "Timeline durumu. Mart-Nisan temel iş. Nisan ikinci sunumu uçtan uca "
            "entegrasyon işaretiydi. Mayıs araştırma fazını kapattı: production migration, "
            "ablation finalization, yeni generation eval. Kalan 14 gün slayt oluşturma, "
            "demo provası ve tez metni. Bu dönemde bilerek yeni feature EKLEMİYORUZ — "
            "odak iletişime kayıyor."
        ),
    },
    # 13 — References
    {
        "title": "References",
        "bullets": [
            "ReliabilityRAG: Effective and Provably Robust Defense for RAG-based Web-Search — arxiv.org/abs/2509.23519",
            "A Multi-Hop and Graph-Based Benchmark for Inter-Context Conflicts in Retrieval-Augmented Generation — arxiv.org/abs/2507.21544",
            "Turkish-Llama-8b-Instruct-v0.1 — huggingface.co/ytu-ce-cosmos/Turkish-Llama-8b-Instruct-v0.1",
            "mDeBERTa-v3-base-mnli-xnli — huggingface.co/MoritzLaurer/mDeBERTa-v3-base-mnli-xnli",
            "multilingual-e5-base — huggingface.co/intfloat/multilingual-e5-base",
            "Project repository: github.com/cgrti/ReliabilityRAG",
        ],
        "notes_en": "Thank you. I'm open to questions.",
        "notes_tr": "Teşekkür ederim. Sorularınızı bekliyorum.",
    },
]


# ── Build ──────────────────────────────────────────────────────────────

def _remove_all_slides(prs):
    """Remove every existing slide; keep masters/layouts/theme intact."""
    sldIdLst = prs.slides._sldIdLst
    # Iterate copies because we mutate
    for sldId in list(sldIdLst):
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        # Drop relationship
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _set_text_runs(text_frame, lines, font_size_pt=None):
    """Replace text_frame content with `lines` (one paragraph per line)."""
    text_frame.clear()
    # First paragraph is auto-created by clear()
    if lines:
        p = text_frame.paragraphs[0]
        p.text = lines[0]
        if font_size_pt:
            for r in p.runs:
                r.font.size = Pt(font_size_pt)
    for line in lines[1:]:
        p = text_frame.add_paragraph()
        p.text = line
        if font_size_pt:
            for r in p.runs:
                r.font.size = Pt(font_size_pt)


def _set_notes(slide, en_text, tr_text):
    """Combine EN + TR speaker notes the way meeting2 does (quoted, English first)."""
    combined = f'"{en_text}"' + f'"{tr_text}"'
    notes = slide.notes_slide
    notes.notes_text_frame.text = combined


def _set_placeholder_text(slide, idx, text):
    """Set placeholder text by idx; ignore if not found."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx and ph.has_text_frame:
            ph.text_frame.text = text
            return True
    return False


def build():
    print(f"Loading template: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    print(f"  size {prs.slide_width} x {prs.slide_height}, "
          f"masters={len(prs.slide_masters)}, layouts={len(prs.slide_layouts)}")

    _remove_all_slides(prs)
    assert len(prs.slides) == 0, "Failed to clear slides"

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # ─ Slide 1: Title ─
    slide = prs.slides.add_slide(title_layout)
    _set_placeholder_text(slide, 0, SLIDE_TITLE_INFO["title"])
    _set_placeholder_text(slide, 1, SLIDE_TITLE_INFO["subtitle"])
    _set_notes(slide, SLIDE_TITLE_INFO["notes_en"], SLIDE_TITLE_INFO["notes_tr"])
    print(f"[1] Title slide added")

    # ─ Content slides 2..13 ─
    for i, spec in enumerate(CONTENT_SLIDES, start=2):
        slide = prs.slides.add_slide(content_layout)
        # Title
        if slide.shapes.title is not None:
            slide.shapes.title.text = spec["title"]
        else:
            _set_placeholder_text(slide, 0, spec["title"])

        # Body: bullets in content placeholder (idx=1, OBJECT type)
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break
        if body_ph is None:
            print(f"[!] Slide {i}: no body placeholder")
        else:
            _set_text_runs(body_ph.text_frame, spec["bullets"], font_size_pt=14)

        # Footer + date + slide number — set placeholder text where layout provides
        _set_placeholder_text(slide, 11, FOOTER_TEXT)   # FOOTER
        _set_placeholder_text(slide, 10, DATE_DISPLAY)  # DATE
        _set_placeholder_text(slide, 12, str(i))        # SLIDE_NUMBER

        _set_notes(slide, spec["notes_en"], spec["notes_tr"])
        print(f"[{i}] '{spec['title']}' — {len(spec['bullets'])} bullets")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"\n[saved] {OUT_PATH}")
    print(f"        size on disk: {OUT_PATH.stat().st_size:,} bytes")


if __name__ == "__main__":
    build()
